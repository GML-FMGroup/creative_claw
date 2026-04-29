"""Native editable PPTX builder for PPT production P0."""

from __future__ import annotations

import html
import zipfile
from pathlib import Path

from src.production.ppt.models import DeckSlide, DeckSpec, PPTRenderSettings
from src.runtime.workspace import workspace_relative_path


class PPTDeckBuilderError(RuntimeError):
    """Raised when native PPTX generation fails."""


class DeckBuilderService:
    """Build editable PPTX files from a deterministic deck specification."""

    def build(self, *, deck_spec: DeckSpec, render_settings: PPTRenderSettings, output_path: Path) -> str:
        """Write an editable PPTX file and return its workspace-relative path."""
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError:
            return _build_minimal_ooxml_deck(
                deck_spec=deck_spec,
                render_settings=render_settings,
                output_path=output_path,
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        theme = _theme_for(render_settings.style_preset)
        prs = Presentation()
        width, height = _slide_size_inches(render_settings.aspect_ratio)
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)
        blank_layout = prs.slide_layouts[6]

        for slide_spec in deck_spec.slides:
            slide = prs.slides.add_slide(blank_layout)
            if slide_spec.layout_type == "cover":
                _add_cover_slide(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt)
            elif slide_spec.layout_type == "closing":
                _add_closing_slide(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt)
            elif slide_spec.layout_type == "metric":
                _add_metric_slide(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt)
            elif slide_spec.layout_type == "two_column":
                _add_two_column_slide(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt)
            else:
                _add_content_slide(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt)
            if slide_spec.speaker_notes:
                try:
                    slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
                except AttributeError:
                    pass

        prs.save(output_path)
        return workspace_relative_path(output_path)

    def build_slide_segments(self, *, deck_spec: DeckSpec, render_settings: PPTRenderSettings, output_dir: Path) -> dict[str, str]:
        """Write one editable single-slide PPTX segment per deck slide."""
        output_dir.mkdir(parents=True, exist_ok=True)
        segment_paths: dict[str, str] = {}
        for slide in deck_spec.slides:
            segment_spec = DeckSpec(
                title=f"{deck_spec.title} - Slide {slide.sequence_index:02d}",
                slides=[slide.model_copy(deep=True)],
            )
            output_path = output_dir / f"slide-{slide.sequence_index:02d}.pptx"
            segment_paths[slide.slide_id] = self.build(
                deck_spec=segment_spec,
                render_settings=render_settings,
                output_path=output_path,
            )
        return segment_paths


def _slide_size_inches(aspect_ratio: str) -> tuple[float, float]:
    if aspect_ratio == "4:3":
        return 10.0, 7.5
    if aspect_ratio == "9:16":
        return 7.5, 13.333
    return 13.333, 7.5


def _theme_for(style_preset: str) -> dict[str, str]:
    themes = {
        "pitch_deck": {
            "dark": "102820",
            "light": "F4EFE7",
            "primary": "2D6A4F",
            "secondary": "95D5B2",
            "accent": "D8A03D",
            "text": "17201A",
        },
        "educational": {
            "dark": "1F2A44",
            "light": "F3F7FB",
            "primary": "2F80ED",
            "secondary": "9CC9F5",
            "accent": "F2994A",
            "text": "172033",
        },
        "editorial_visual": {
            "dark": "2B1B17",
            "light": "F7EDE2",
            "primary": "B85042",
            "secondary": "A7BEAE",
            "accent": "F2C078",
            "text": "241713",
        },
        "business_executive": {
            "dark": "1E2761",
            "light": "F6F8FE",
            "primary": "334EAC",
            "secondary": "CADCFC",
            "accent": "F9B115",
            "text": "182033",
        },
    }
    return themes.get(style_preset, themes["business_executive"])


def _set_bg(slide, color: str, RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _text_box(slide, text: str, *, x: float, y: float, w: float, h: float, font_size: int, color: str, Inches, Pt, RGBColor, bold: bool = False, font_face: str = "Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.name = font_face
    paragraph.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_bullets(slide, bullets: list[str], *, x: float, y: float, w: float, h: float, theme: dict[str, str], Inches, Pt, RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets[:6]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18 if len(bullets) <= 4 else 15)
        paragraph.font.color.rgb = RGBColor.from_string(theme["text"])
        paragraph.space_after = Pt(8)


def _add_footer(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, Inches, Pt) -> None:
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(height - 0.32),
        Inches(width),
        Inches(0.32),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor.from_string(theme["dark"])
    footer.line.fill.background()
    label = f"{slide_spec.sequence_index:02d}"
    _text_box(slide, label, x=width - 0.75, y=height - 0.25, w=0.45, h=0.15, font_size=9, color=theme["secondary"], Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _add_cover_slide(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt) -> None:
    _set_bg(slide, theme["dark"], RGBColor)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.55), Inches(height))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])
    accent.line.fill.background()
    _text_box(slide, slide_spec.title, x=0.95, y=1.45, w=width - 1.9, h=1.15, font_size=42, color="FFFFFF", bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    subtitle = slide_spec.bullets[0] if slide_spec.bullets else slide_spec.visual_notes
    _text_box(slide, subtitle, x=0.98, y=3.0, w=width - 2.0, h=1.3, font_size=20, color=theme["secondary"], Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    tag = _text_box(slide, "CREATIVECLAW PPT PRODUCTION", x=0.98, y=height - 1.0, w=4.2, h=0.25, font_size=10, color=theme["accent"], bold=True, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _add_content_slide(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt) -> None:
    _set_bg(slide, theme["light"], RGBColor)
    _text_box(slide, slide_spec.title, x=0.65, y=0.48, w=width - 1.3, h=0.58, font_size=30, color=theme["dark"], bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.45), Inches(width - 1.3), Inches(height - 2.25))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    card.line.color.rgb = RGBColor.from_string(theme["secondary"])
    _add_bullets(slide, slide_spec.bullets, x=1.05, y=1.86, w=width - 2.1, h=height - 3.0, theme=theme, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _add_footer(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, Inches, Pt)


def _add_two_column_slide(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt) -> None:
    _set_bg(slide, theme["light"], RGBColor)
    _text_box(slide, slide_spec.title, x=0.65, y=0.48, w=width - 1.3, h=0.58, font_size=30, color=theme["dark"], bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.45), Inches((width - 1.6) / 2), Inches(height - 2.25))
    left.fill.solid(); left.fill.fore_color.rgb = RGBColor.from_string("FFFFFF"); left.line.color.rgb = RGBColor.from_string(theme["secondary"])
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(width / 2 + 0.15), Inches(1.45), Inches((width - 1.6) / 2), Inches(height - 2.25))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor.from_string(theme["secondary"]); right.line.fill.background()
    midpoint = max(1, (len(slide_spec.bullets) + 1) // 2)
    _add_bullets(slide, slide_spec.bullets[:midpoint], x=0.95, y=1.85, w=(width - 2.2) / 2, h=height - 3.0, theme=theme, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _add_bullets(slide, slide_spec.bullets[midpoint:] or [slide_spec.visual_notes or "Use this space for a visual proof point."], x=width / 2 + 0.45, y=1.85, w=(width - 2.2) / 2, h=height - 3.0, theme=theme, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _add_footer(slide, slide_spec, theme, width, height, RGBColor, MSO_SHAPE, Inches, Pt)


def _add_metric_slide(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt) -> None:
    _set_bg(slide, theme["dark"], RGBColor)
    _text_box(slide, slide_spec.title, x=0.7, y=0.55, w=width - 1.4, h=0.6, font_size=30, color="FFFFFF", bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    bullets = slide_spec.bullets[:3] or ["Key metric", "Business implication", "Next action"]
    card_w = (width - 1.7) / 3
    for index, bullet in enumerate(bullets):
        x = 0.7 + index * (card_w + 0.15)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(card_w), Inches(height - 2.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string(theme["light"]); card.line.fill.background()
        _text_box(slide, f"0{index + 1}", x=x + 0.25, y=2.08, w=card_w - 0.5, h=0.6, font_size=34, color=theme["accent"], bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        _text_box(slide, bullet, x=x + 0.25, y=3.05, w=card_w - 0.5, h=1.6, font_size=16, color=theme["text"], bold=True, Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _add_closing_slide(slide, slide_spec: DeckSlide, theme: dict[str, str], width: float, height: float, RGBColor, MSO_SHAPE, PP_ALIGN, Inches, Pt) -> None:
    _set_bg(slide, theme["dark"], RGBColor)
    _text_box(slide, slide_spec.title, x=0.95, y=1.35, w=width - 1.9, h=0.95, font_size=40, color="FFFFFF", bold=True, font_face="Georgia", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _add_bullets(slide, slide_spec.bullets or ["Confirm decisions", "Assign owners", "Move to execution"], x=1.0, y=2.8, w=width - 2.0, h=2.2, theme={**theme, "text": "FFFFFF"}, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(width - 1.4), Inches(height - 1.35), Inches(0.72), Inches(0.72))
    mark.fill.solid(); mark.fill.fore_color.rgb = RGBColor.from_string(theme["accent"]); mark.line.fill.background()


def _build_minimal_ooxml_deck(*, deck_spec: DeckSpec, render_settings: PPTRenderSettings, output_path: Path) -> str:
    """Build a minimal editable PPTX package without third-party dependencies."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    width, height = _slide_size_emu(render_settings.aspect_ratio)
    theme = _theme_for(render_settings.style_preset)
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", _content_types(len(deck_spec.slides)))
        package.writestr("_rels/.rels", _root_rels())
        package.writestr("docProps/app.xml", _app_xml(len(deck_spec.slides)))
        package.writestr("docProps/core.xml", _core_xml(deck_spec.title))
        package.writestr("ppt/presentation.xml", _presentation_xml(len(deck_spec.slides), width, height))
        package.writestr("ppt/_rels/presentation.xml.rels", _presentation_rels(len(deck_spec.slides)))
        package.writestr("ppt/theme/theme1.xml", _theme_xml())
        package.writestr("ppt/slideMasters/slideMaster1.xml", _slide_master_xml())
        package.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", _slide_master_rels())
        package.writestr("ppt/slideLayouts/slideLayout1.xml", _slide_layout_xml())
        package.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", _slide_layout_rels())
        for part_index, slide in enumerate(deck_spec.slides, start=1):
            package.writestr(f"ppt/slides/slide{part_index}.xml", _slide_xml(slide, theme))
            package.writestr(f"ppt/slides/_rels/slide{part_index}.xml.rels", _slide_rels())
    return workspace_relative_path(output_path)


def _slide_size_emu(aspect_ratio: str) -> tuple[int, int]:
    if aspect_ratio == "4:3":
        return 9144000, 6858000
    if aspect_ratio == "9:16":
        return 6858000, 12192000
    return 12192000, 6858000


def _content_types(slide_count: int) -> str:
    slide_overrides = "".join(
        f'<Override PartName="/ppt/slides/slide{index}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for index in range(1, slide_count + 1)
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  {slide_overrides}
</Types>'''


def _root_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>'''


def _app_xml(slide_count: int) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>CreativeClaw</Application>
  <PresentationFormat>On-screen Show</PresentationFormat>
  <Slides>{slide_count}</Slides>
</Properties>'''


def _core_xml(title: str) -> str:
    escaped = _xml_text(title or "CreativeClaw Deck")
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>{escaped}</dc:title>
  <dc:creator>CreativeClaw</dc:creator>
</cp:coreProperties>'''


def _presentation_xml(slide_count: int, width: int, height: int) -> str:
    slide_ids = "".join(f'<p:sldId id="{255 + index}" r:id="rId{index + 1}"/>' for index in range(1, slide_count + 1))
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slide_ids}</p:sldIdLst>
  <p:sldSz cx="{width}" cy="{height}"/>
  <p:notesSz cx="6858000" cy="9144000"/>
  <p:defaultTextStyle/>
</p:presentation>'''


def _presentation_rels(slide_count: int) -> str:
    relationships = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>']
    relationships.extend(
        f'<Relationship Id="rId{index + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{index}.xml"/>'
        for index in range(1, slide_count + 1)
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">{''.join(relationships)}</Relationships>'''


def _theme_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="CreativeClaw">
  <a:themeElements>
    <a:clrScheme name="CreativeClaw"><a:dk1><a:srgbClr val="111111"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="1E2761"/></a:dk2><a:lt2><a:srgbClr val="F6F8FE"/></a:lt2><a:accent1><a:srgbClr val="334EAC"/></a:accent1><a:accent2><a:srgbClr val="F9B115"/></a:accent2><a:accent3><a:srgbClr val="CADCFC"/></a:accent3><a:accent4><a:srgbClr val="2D6A4F"/></a:accent4><a:accent5><a:srgbClr val="B85042"/></a:accent5><a:accent6><a:srgbClr val="A7BEAE"/></a:accent6><a:hlink><a:srgbClr val="0000FF"/></a:hlink><a:folHlink><a:srgbClr val="800080"/></a:folHlink></a:clrScheme>
    <a:fontScheme name="CreativeClaw"><a:majorFont><a:latin typeface="Georgia"/></a:majorFont><a:minorFont><a:latin typeface="Aptos"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="CreativeClaw"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme>
  </a:themeElements>
</a:theme>'''


def _slide_master_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
  <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>
</p:sldMaster>'''


def _slide_master_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>'''


def _slide_layout_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">
  <p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>'''


def _slide_layout_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>'''


def _slide_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>'''


def _slide_xml(slide: DeckSlide, theme: dict[str, str]) -> str:
    background = theme["dark"] if slide.layout_type in {"cover", "closing", "metric"} else theme["light"]
    title_color = "FFFFFF" if slide.layout_type in {"cover", "closing", "metric"} else theme["dark"]
    body_color = "FFFFFF" if slide.layout_type in {"cover", "closing", "metric"} else theme["text"]
    shapes = [
        _text_shape(2, "Title", slide.title, 520000, 520000, 10800000, 900000, 3400, title_color, bold=True, font="Georgia"),
    ]
    y = 1550000
    for idx, bullet in enumerate((slide.bullets or [slide.visual_notes])[:7], start=3):
        shapes.append(_text_shape(idx, f"Bullet {idx}", bullet, 760000, y, 10300000, 520000, 1800, body_color, bold=False, font="Aptos"))
        y += 620000
    shapes.append(_text_shape(40, "Slide Number", f"{slide.sequence_index:02d}", 11100000, 6250000, 700000, 260000, 1100, title_color, bold=False, font="Aptos"))
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val="{background}"/></a:solidFill><a:effectLst/></p:bgPr></p:bg><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>{''.join(shapes)}</p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''


def _text_shape(shape_id: int, name: str, text: str, x: int, y: int, w: int, h: int, size: int, color: str, *, bold: bool, font: str) -> str:
    bold_attr = ' b="1"' if bold else ""
    return f'''<p:sp><p:nvSpPr><p:cNvPr id="{shape_id}" name="{_xml_attr(name)}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr><p:txBody><a:bodyPr wrap="square"><a:spAutoFit/></a:bodyPr><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="{size}"{bold_attr}><a:solidFill><a:srgbClr val="{color}"/></a:solidFill><a:latin typeface="{_xml_attr(font)}"/></a:rPr><a:t>{_xml_text(text)}</a:t></a:r><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>'''


def _xml_text(value: str) -> str:
    return html.escape(str(value or ""), quote=False)


def _xml_attr(value: str) -> str:
    return html.escape(str(value or ""), quote=True)
