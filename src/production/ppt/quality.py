"""Deterministic quality reporting for PPT production."""

from __future__ import annotations

from dataclasses import dataclass
import re
import zipfile
from pathlib import Path

from src.production.ppt.models import PPTProductionState, PPTQualityCheck, PPTQualityReport
from src.runtime.workspace import resolve_workspace_path


@dataclass(frozen=True)
class _SlideCountInspection:
    """Result of inspecting generated PPTX slide count."""

    count: int
    method: str
    error: str = ""


def build_quality_report(state: PPTProductionState, *, report_path: str | None = None) -> PPTQualityReport:
    """Build an explainable quality report from persisted PPT production state."""
    checks = [
        _final_pptx_exists_check(state),
        _slide_count_check(state),
        _outline_check(state),
        _preview_check(state),
        _input_warning_check(state),
        _content_density_check(state),
        _source_fact_coverage_check(state),
    ]
    status = _aggregate_status(checks)
    recommendations = _recommendations(checks)
    metrics = _metrics(state)
    return PPTQualityReport(
        status=status,
        summary=_summary_for(status, checks),
        report_path=report_path,
        metrics=metrics,
        checks=checks,
        recommendations=recommendations,
    )


def quality_report_markdown(report: PPTQualityReport | None) -> str:
    """Render a PPT quality report as operator-readable Markdown."""
    if report is None:
        return "# PPT Quality Report\n\nNo quality report has been generated yet.\n"
    lines = [
        "# PPT Quality Report",
        "",
        f"- Status: {report.status}",
        f"- Summary: {report.summary}",
        "",
        "## Metrics",
        "",
    ]
    if report.metrics:
        lines.extend(f"- {key}: {value}" for key, value in report.metrics.items())
    else:
        lines.append("- No metrics available.")
    lines.extend(["", "## Checks", ""])
    if report.checks:
        for check in report.checks:
            lines.extend(_check_markdown_lines(check))
    else:
        lines.append("- No checks available.")
    lines.extend(["", "## Recommendations", ""])
    if report.recommendations:
        lines.extend(f"- {item}" for item in report.recommendations)
    else:
        lines.append("- No immediate quality actions found.")
    return "\n".join(lines).rstrip() + "\n"


def _check_markdown_lines(check: PPTQualityCheck) -> list[str]:
    lines = [
        f"### [{check.status.upper()}] {check.check_id}",
        "",
        f"- Category: {check.category}",
        f"- Summary: {check.summary}",
    ]
    detail_lines = _detail_markdown_lines(check.details)
    if detail_lines:
        lines.extend(["- Details:", *detail_lines])
    else:
        lines.append("- Details: none")
    lines.append("")
    return lines


def _detail_markdown_lines(details: dict[str, object]) -> list[str]:
    lines: list[str] = []
    for key, value in details.items():
        lines.extend(_detail_value_lines(str(key), value, indent=2))
    return lines


def _detail_value_lines(label: str, value: object, *, indent: int) -> list[str]:
    prefix = " " * indent
    if isinstance(value, dict):
        if not value:
            return [f"{prefix}- {label}: {{}}"]
        lines = [f"{prefix}- {label}:"]
        for nested_key, nested_value in value.items():
            lines.extend(_detail_value_lines(str(nested_key), nested_value, indent=indent + 2))
        return lines
    if isinstance(value, list):
        if not value:
            return [f"{prefix}- {label}: []"]
        lines = [f"{prefix}- {label}:"]
        item_prefix = " " * (indent + 2)
        for item in value:
            if isinstance(item, dict):
                lines.append(f"{item_prefix}-")
                for nested_key, nested_value in item.items():
                    lines.extend(_detail_value_lines(str(nested_key), nested_value, indent=indent + 4))
            else:
                lines.append(f"{item_prefix}- {_format_detail_scalar(item)}")
        return lines
    return [f"{prefix}- {label}: {_format_detail_scalar(value)}"]


def _format_detail_scalar(value: object) -> str:
    if value is None:
        return "null"
    if isinstance(value, bool):
        return "true" if value else "false"
    return str(value)


def _final_pptx_exists_check(state: PPTProductionState) -> PPTQualityCheck:
    path = state.final_artifact.pptx_path if state.final_artifact is not None else ""
    exists = bool(path and resolve_workspace_path(path).is_file()) if path else False
    return _check(
        "final_pptx_exists",
        "structure",
        "pass" if exists else "fail",
        "Final PPTX file exists." if exists else "No final PPTX file is recorded or present on disk.",
        {"path": path},
    )


def _slide_count_check(state: PPTProductionState) -> PPTQualityCheck:
    expected = len(state.deck_spec.slides) if state.deck_spec is not None else 0
    inspection = _inspect_pptx_slide_count(state)
    actual = inspection.count
    details: dict[str, object] = {
        "expected": expected,
        "actual": actual,
        "inspection_method": inspection.method,
    }
    if inspection.error:
        details["inspection_error"] = inspection.error
    if expected <= 0:
        return _check("slide_count", "structure", "fail", "No deck spec slides are available.")
    if actual <= 0:
        return _check("slide_count", "structure", "warning", "Could not inspect generated PPTX slide count.", details)
    return _check(
        "slide_count",
        "structure",
        "pass" if actual == expected else "warning",
        "Generated PPTX slide count matches the deck spec." if actual == expected else "Generated PPTX slide count differs from the deck spec.",
        details,
    )


def _outline_check(state: PPTProductionState) -> PPTQualityCheck:
    if state.outline is None or not state.outline.entries:
        return _check("outline_present", "content", "fail", "No outline exists for this PPT production.")
    empty_titles = [entry.sequence_index for entry in state.outline.entries if not entry.title.strip()]
    return _check(
        "outline_present",
        "content",
        "pass" if not empty_titles else "warning",
        "Outline exists and all slides have titles." if not empty_titles else "Some outline entries are missing titles.",
        {"empty_title_slides": empty_titles},
    )


def _preview_check(state: PPTProductionState) -> PPTQualityCheck:
    expected = len(state.deck_spec.slides) if state.deck_spec is not None else 0
    valid = [item for item in state.slide_previews if item.status == "generated" and resolve_workspace_path(item.preview_path).is_file()]
    return _check(
        "preview_images",
        "visual",
        "pass" if expected and len(valid) >= expected else "warning",
        "Preview images were generated for all slides." if expected and len(valid) >= expected else "Preview images are missing or incomplete.",
        {"expected": expected, "actual": len(valid)},
    )


def _input_warning_check(state: PPTProductionState) -> PPTQualityCheck:
    warnings = [item.warning for item in state.inputs if item.warning]
    warnings.extend(state.warnings)
    return _check(
        "input_support",
        "delivery",
        "warning" if warnings else "pass",
        "Some inputs were recorded but not fully used in P0." if warnings else "No input support warnings were recorded.",
        {"warnings": warnings},
    )


def _content_density_check(state: PPTProductionState) -> PPTQualityCheck:
    if state.deck_spec is None:
        return _check("content_density", "content", "not_applicable", "No deck spec exists.")
    dense_slides = [slide.sequence_index for slide in state.deck_spec.slides if len(" ".join(slide.bullets)) > 620 or len(slide.bullets) > 6]
    return _check(
        "content_density",
        "content",
        "warning" if dense_slides else "pass",
        "Some slides may be too dense for presentation use." if dense_slides else "Slide text density is within P0 limits.",
        {"dense_slides": dense_slides},
    )


def _source_fact_coverage_check(state: PPTProductionState) -> PPTQualityCheck:
    """Check whether extracted source-document facts are represented in the deck."""
    document_summary = state.document_summary
    if document_summary is None or document_summary.status != "ready":
        return _check(
            "source_fact_coverage",
            "content",
            "not_applicable",
            "No ready source-document facts are available for coverage checking.",
        )
    facts = [fact.strip() for fact in document_summary.salient_facts if fact.strip()]
    if not facts:
        return _check(
            "source_fact_coverage",
            "content",
            "not_applicable",
            "No salient source-document facts are available for coverage checking.",
        )

    deck_text = _combined_deck_text(state)
    matched_facts = [fact for fact in facts if _fact_matches_text(fact, deck_text)]
    matched_count = len(matched_facts)
    coverage_ratio = round(matched_count / len(facts), 3)
    details = {
        "source_input_ids": document_summary.source_input_ids,
        "source_document_details": _source_document_details(state, document_summary.source_input_ids),
        "fact_count": len(facts),
        "matched_fact_count": matched_count,
        "coverage_ratio": coverage_ratio,
        "matched_facts": matched_facts[:3],
        "unmatched_facts": [fact for fact in facts if fact not in matched_facts][:3],
    }
    if matched_count:
        return _check(
            "source_fact_coverage",
            "content",
            "pass",
            "At least one extracted source-document fact is represented in the deck.",
            details,
        )
    return _check(
        "source_fact_coverage",
        "content",
        "warning",
        "Extracted source-document facts were not found in the generated deck content.",
        details,
    )


def _combined_deck_text(state: PPTProductionState) -> str:
    """Return the reviewable outline/deck text used for source fact coverage."""
    parts: list[str] = []
    if state.outline is not None:
        parts.append(state.outline.title)
        for entry in state.outline.entries:
            parts.extend(
                [
                    entry.title,
                    entry.purpose,
                    entry.speaker_notes,
                    *entry.bullet_points,
                ]
            )
    if state.deck_spec is not None:
        parts.append(state.deck_spec.title)
        for slide in state.deck_spec.slides:
            parts.extend(
                [
                    slide.title,
                    slide.visual_notes,
                    slide.speaker_notes,
                    *slide.bullets,
                ]
            )
    return "\n".join(part for part in parts if part)


def _source_document_details(state: PPTProductionState, source_input_ids: list[str]) -> list[dict[str, str]]:
    """Resolve source document input ids to readable metadata for quality reports."""
    inputs_by_id = {item.input_id: item for item in state.inputs}
    details: list[dict[str, str]] = []
    for input_id in source_input_ids:
        entry = inputs_by_id.get(input_id)
        if entry is None:
            details.append({"input_id": input_id, "name": input_id, "path": "", "role": "", "status": "missing"})
            continue
        details.append(
            {
                "input_id": entry.input_id,
                "name": entry.name,
                "path": entry.path,
                "role": entry.role,
                "status": entry.status,
            }
        )
    return details


def _fact_matches_text(fact: str, text: str) -> bool:
    normalized_fact = _normalize_match_text(fact)
    normalized_text = _normalize_match_text(text)
    if not normalized_fact or not normalized_text:
        return False
    if normalized_fact in normalized_text:
        return True

    fact_tokens = set(normalized_fact.split())
    if len(fact_tokens) < 4:
        return fact_tokens.issubset(set(normalized_text.split()))
    matched = fact_tokens.intersection(normalized_text.split())
    return len(matched) / len(fact_tokens) >= 0.65


def _normalize_match_text(text: str) -> str:
    tokens = re.findall(r"[a-z0-9]+|[\u4e00-\u9fff]", text.lower())
    return " ".join(tokens)


def _inspect_pptx_slide_count(state: PPTProductionState) -> _SlideCountInspection:
    """Inspect final PPTX slide count and preserve diagnostic failure details."""
    path = state.final_artifact.pptx_path if state.final_artifact is not None else ""
    if not path:
        return _SlideCountInspection(count=0, method="missing_path", error="No final PPTX path is recorded.")
    try:
        from pptx import Presentation

        return _SlideCountInspection(count=len(Presentation(str(resolve_workspace_path(path))).slides), method="python_pptx")
    except Exception as pptx_exc:
        zip_count, zip_error = _zip_slide_count(path)
        pptx_error = f"python_pptx_failed:{type(pptx_exc).__name__}: {pptx_exc}"
        if zip_count > 0:
            return _SlideCountInspection(count=zip_count, method="zip_package", error=pptx_error)
        return _SlideCountInspection(count=0, method="failed", error=f"{pptx_error}; zip_failed:{zip_error}")


def _zip_slide_count(path: str) -> tuple[int, str]:
    try:
        with zipfile.ZipFile(resolve_workspace_path(path)) as package:
            slide_count = len(
                [
                    name
                    for name in package.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                ]
            )
            return slide_count, "" if slide_count else "no_slide_xml_found"
    except Exception as exc:
        return 0, f"{type(exc).__name__}: {exc}"


def _metrics(state: PPTProductionState) -> dict[str, object]:
    return {
        "target_pages": state.render_settings.target_pages,
        "outline_slides": len(state.outline.entries) if state.outline is not None else 0,
        "deck_spec_slides": len(state.deck_spec.slides) if state.deck_spec is not None else 0,
        "preview_images": len(state.slide_previews),
        "input_count": len(state.inputs),
        "warnings": len(state.warnings) + len([item for item in state.inputs if item.warning]),
    }


def _aggregate_status(checks: list[PPTQualityCheck]) -> str:
    if any(check.status == "fail" for check in checks):
        return "fail"
    if any(check.status == "warning" for check in checks):
        return "warning"
    return "pass"


def _summary_for(status: str, checks: list[PPTQualityCheck]) -> str:
    failing = len([check for check in checks if check.status == "fail"])
    warnings = len([check for check in checks if check.status == "warning"])
    if status == "fail":
        return f"PPT generation completed with {failing} failing check(s) and {warnings} warning(s)."
    if status == "warning":
        return f"PPT generation completed with {warnings} warning(s)."
    return "PPT generation completed and deterministic checks passed."


def _recommendations(checks: list[PPTQualityCheck]) -> list[str]:
    recommendations: list[str] = []
    for check in checks:
        if check.status == "pass" or check.status == "not_applicable":
            continue
        if check.check_id == "input_support":
            recommendations.append("If template or source-document fidelity matters, continue with the P1 template/document pipeline.")
        elif check.check_id == "content_density":
            recommendations.append("Split dense slides or shorten bullets before final delivery.")
        elif check.check_id == "preview_images":
            recommendations.append("Install or verify LibreOffice/Poppler if real slide previews are required.")
        else:
            recommendations.append(f"Review check `{check.check_id}` before delivery.")
    return recommendations


def _check(check_id: str, category: str, status: str, summary: str, details: dict[str, object] | None = None) -> PPTQualityCheck:
    return PPTQualityCheck(
        check_id=check_id,
        category=category,  # type: ignore[arg-type]
        status=status,  # type: ignore[arg-type]
        summary=summary,
        details=details or {},
    )
